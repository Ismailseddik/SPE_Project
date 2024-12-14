import socket
import os
import threading
import random
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
import time

HOST = '127.0.0.1'
PORT = 65432

# Function to generate a dummy Word (.docx) file with a heading and 50 paragraphs
def generate_docx(file_name):
    """Generate a dummy Word (.docx) file."""
    doc = Document()
    doc.add_heading("Test Document", level=1)
    for i in range(50):  # Add 50 dummy paragraphs
        doc.add_paragraph(f"This is paragraph {i}.")
    doc.save(file_name)

# Function to generate a dummy Excel (.xlsx) file with 50 rows and 10 columns of data
def generate_xlsx(file_name):
    """Generate a dummy Excel (.xlsx) file."""
    wb = Workbook()
    ws = wb.active
    for i in range(50):  # Add 50 rows of dummy data
        ws.append([f"Data {i}-{j}" for j in range(10)])  # 10 columns per row
    wb.save(file_name)

# Function to generate a dummy PowerPoint (.pptx) file with 10 slides
def generate_pptx(file_name):
    """Generate a dummy PowerPoint (.pptx) file."""
    prs = Presentation()
    for i in range(10):  # Add 10 dummy slides
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide layout
        slide.shapes.title.text = f"Slide {i}"
    prs.save(file_name)

# Function to generate a file of the specified type (docx, xlsx, pptx)
def generate_dummy_file(file_type, file_name):
    """Generate a file of the specified type."""
    if file_type == "docx":
        generate_docx(file_name)
    elif file_type == "xlsx":
        generate_xlsx(file_name)
    elif file_type == "pptx":
        generate_pptx(file_name)

# Function to send a file to the server and handle the server's response
def send_file(file_name):
    """Send a file to the server and wait for the response."""
    file_size = os.path.getsize(file_name)
    with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as client:
        try:
            client.connect((HOST, PORT))
            
            # Step 1: Send Metadata (file name and size)
            metadata = f"{file_name};{file_size}"
            client.send(metadata.encode('utf-8'))

            # Step 2: Wait for Metadata Acknowledgment from the server
            ack = client.recv(1024).decode('utf-8')
            if ack != "metadata_received":
                print(f"[ERROR] Metadata acknowledgment failed for {file_name}")
                return

            # Step 3: Send File Data in chunks of 1024 bytes
            with open(file_name, 'rb') as f:
                while (chunk := f.read(1024)):
                    client.sendall(chunk)
            print(f"[INFO] File {file_name} sent successfully.")

            # Step 4: Wait for Server Response regarding file conversion
            response = client.recv(1024).decode('utf-8')
            if response == "conversion_successful":
                print(f"[INFO] File {file_name} converted successfully.")
            elif response == "unsupported_file_type":
                print(f"[ERROR] Unsupported file type for {file_name}.")
            elif response == "file_transfer_error":
                print(f"[ERROR] File transfer error for {file_name}.")
            elif response == "conversion_failed":
                print(f"[ERROR] Conversion failed for {file_name}.")
            else:
                print(f"[ERROR] Unknown server response for {file_name}: {response}")
        except Exception as e:
            print(f"[ERROR] Exception occurred while sending {file_name}: {e}")

# Function to perform a stress test by sending multiple large files concurrently
def stress_test(num_files, file_size_kb):
    """Perform a stress test by sending multiple large files."""
    threads = []

    for i in range(num_files):
        # Randomly select a file type from docx, xlsx, pptx
        file_type = random.choice(["docx", "xlsx", "pptx"])
        file_name = f"dummy_file_{i}.{file_type}"

        # Generate a dummy file of the selected type
        generate_dummy_file(file_type, file_name)

        # Create a thread to send the file to the server
        t = threading.Thread(target=send_file, args=(file_name,))
        threads.append(t)
        t.start()

    # Wait for all threads to complete
    for t in threads:
        t.join()

if __name__ == "__main__":
    print("Starting stress test...")
    start_time = time.time()

    # Stress test with 10 files, each approximately 1MB
    stress_test(num_files=10, file_size_kb=1024)

    end_time = time.time()
    print(f"Stress test completed in {end_time - start_time:.2f} seconds.")
