from pptx import Presentation 
from fpdf import FPDF
def convert_ppt_to_pdf(file_name): 
    try:
        prs = Presentation(file_name) 
		pdf = FPDF()
		pdf.add_page() 
		pdf.set_font("Arial", size=12) 
		for slide in prs.slides:
	    	for shape in slide.shapes: 
				if shape.has_text_frame:
		   			for para in shape.text_frame.paragraphs: 			
            			pdf.multi_cell(0, 10, para.text)
		output_file = file_name.replace('.pptx', '.pdf') 	
  		pdf.output(output_file)
		return output_file 
    except Exception as e:
		print(f"Error converting PowerPoint to PDF: {e}") 	
  		return None